import os
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from PIL import Image, ImageDraw, ImageFont

# -----------------------------------------------------------------------------
# 1. إعداد المسارات والمجلدات
# -----------------------------------------------------------------------------
OUTPUT_DIR = "output"
IMAGES_DIR = "images"
os.makedirs(OUTPUT_DIR, exist_ok=True)
os.makedirs(IMAGES_DIR, exist_ok=True)

IMG_HIBISCUS = os.path.join(IMAGES_DIR, "hibiscus.png")
IMG_LILY = os.path.join(IMAGES_DIR, "lily.png")
IMG_PEONY = os.path.join(IMAGES_DIR, "peony.png")

# -----------------------------------------------------------------------------
# 2. توليد صور افتراضية عالية الدقة في حال لم يضع المستخدم صوره بعد
# -----------------------------------------------------------------------------
def ensure_images():
    if not os.path.exists(IMG_HIBISCUS):
        img = Image.new("RGBA", (600, 600), (255, 255, 255, 0))
        d = ImageDraw.Draw(img)
        d.ellipse([50, 50, 550, 550], fill=(255, 210, 60, 255))
        d.ellipse([220, 220, 380, 380], fill=(180, 40, 20, 255))
        img.save(IMG_HIBISCUS)
        
    if not os.path.exists(IMG_LILY):
        img = Image.new("RGBA", (600, 600), (255, 255, 255, 0))
        d = ImageDraw.Draw(img)
        d.ellipse([70, 70, 530, 530], fill=(245, 140, 175, 255))
        d.ellipse([250, 250, 350, 350], fill=(255, 255, 200, 255))
        img.save(IMG_LILY)

    if not os.path.exists(IMG_PEONY):
        img = Image.new("RGBA", (600, 600), (255, 255, 255, 0))
        d = ImageDraw.Draw(img)
        d.ellipse([60, 60, 540, 540], fill=(70, 130, 180, 255))
        d.ellipse([230, 230, 370, 370], fill=(30, 60, 100, 255))
        img.save(IMG_PEONY)

ensure_images()

# -----------------------------------------------------------------------------
# 3. الألوان المحددة لكل بطاقة (مطابقة للمقطع)
# -----------------------------------------------------------------------------
COLOR_BG_1 = RGBColor(253, 244, 227)   # الكركديه (أصفر دافئ)
COLOR_BG_2 = RGBColor(252, 228, 236)   # الزنبق (وردي باستيل)
COLOR_BG_3 = RGBColor(227, 238, 246)   # الفاوانيا (أزرق هادئ)

ACCENT_1 = RGBColor(218, 120, 23)
ACCENT_2 = RGBColor(216, 67, 108)
ACCENT_3 = RGBColor(35, 78, 112)

TEXT_1 = RGBColor(90, 65, 40)
TEXT_2 = RGBColor(95, 45, 65)
TEXT_3 = RGBColor(40, 60, 80)

# -----------------------------------------------------------------------------
# 4. بناء ملف PowerPoint مع وسم الـ Morph المتقدم
# -----------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

def apply_precise_morph(slide, duration_ms=1250):
    # استخدام وسم الـ Morph الخاص بمايكروسوفت مع تحديد المدة بدقة
    xml = f"""
    <p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" 
                  xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" 
                  spd="med" advClick="1">
      <p14:morph option="byObject"/>
    </p:transition>
    """
    slide._element.append(parse_xml(xml))

def create_powerpoint_slide(state_index):
    slide = prs.slides.add_slide(blank_layout)
    apply_precise_morph(slide)

    # حسابات الحالات (State Definition)
    # الحالة 0: وضع التساوي
    # الحالة 1: توسيع البطاقة الأولى
    # الحالة 2: توسيع البطاقة الثانية
    # الحالة 3: توسيع البطاقة الثالثة
    if state_index == 0:
        w = [Inches(4.444), Inches(4.444), Inches(4.445)]
        x = [Inches(0.0), Inches(4.444), Inches(8.888)]
    elif state_index == 1:
        w = [Inches(9.133), Inches(2.100), Inches(2.100)]
        x = [Inches(0.0), Inches(9.133), Inches(11.233)]
    elif state_index == 2:
        w = [Inches(2.100), Inches(9.133), Inches(2.100)]
        x = [Inches(0.0), Inches(2.100), Inches(11.233)]
    elif state_index == 3:
        w = [Inches(2.100), Inches(2.100), Inches(9.133)]
        x = [Inches(0.0), Inches(2.100), Inches(4.200)]

    # رسم المستطيلات الأساسية مع أسماء الـ Morph الثابتة (!!)
    bg1 = slide.shapes.add_shape(1, x[0], Inches(0), w[0], Inches(7.5))
    bg1.name = "!!card_1"
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_BG_1
    bg1.line.fill.background()

    bg2 = slide.shapes.add_shape(1, x[1], Inches(0), w[1], Inches(7.5))
    bg2.name = "!!card_2"
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = COLOR_BG_2
    bg2.line.fill.background()

    bg3 = slide.shapes.add_shape(1, x[2], Inches(0), w[2], Inches(7.5))
    bg3.name = "!!card_3"
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = COLOR_BG_3
    bg3.line.fill.background()

    # إضافة العناصر لكل بطاقة حسب حالتها
    # --- بطاقة 1 ---
    if state_index == 1:
        f1 = slide.shapes.add_picture(IMG_HIBISCUS, Inches(0.4), Inches(1.2), Inches(4.7), Inches(4.7))
        f1.name = "!!flower_1"
        
        num1 = slide.shapes.add_textbox(Inches(5.0), Inches(0.8), Inches(3.5), Inches(1.8))
        num1.name = "!!num_1"
        p = num1.text_frame.paragraphs[0]
        p.text = "1"
        p.font.name = "Segoe UI"
        p.font.size = Pt(110)
        p.font.bold = True
        p.font.color.rgb = ACCENT_1

        t1 = slide.shapes.add_textbox(Inches(4.7), Inches(2.8), Inches(4.2), Inches(3.6))
        t1.name = "!!text_1"
        tf = t1.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = "رمز للإشراق والبهجة والجمال الاستوائي، تمنح الأجواء طاقة إيجابية وحيوية متجددة."
        p1.font.name = "Calibri"
        p1.font.size = Pt(18)
        p1.font.color.rgb = TEXT_1
        p1.alignment = PP_ALIGN.RIGHT
        p1.space_after = Pt(18)
        p2 = tf.add_paragraph()
        p2.text = "تتميز ببتلاتها المتوهجة للشمس، وتعد خياراً مثالياً للاحتفاء بالبدايات السعيدة ولحظات الفرح."
        p2.font.name = "Calibri"
        p2.font.size = Pt(18)
        p2.font.color.rgb = TEXT_1
        p2.alignment = PP_ALIGN.RIGHT

        lbl1 = slide.shapes.add_textbox(Inches(0.7), Inches(5.9), Inches(4.1), Inches(1.2))
        lbl1.name = "!!label_1"
        lbl1.text_frame.paragraphs[0].text = "Yellow Hibiscus"
        lbl1.text_frame.paragraphs[0].font.name = "Georgia"
        lbl1.text_frame.paragraphs[0].font.size = Pt(24)
        lbl1.text_frame.paragraphs[0].font.italic = True
        lbl1.text_frame.paragraphs[0].font.color.rgb = ACCENT_1
        lbl1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        p_ar = lbl1.text_frame.add_paragraph()
        p_ar.text = "الكركديه الأصفر"
        p_ar.font.name = "Calibri"
        p_ar.font.size = Pt(16)
        p_ar.font.color.rgb = TEXT_1
        p_ar.alignment = PP_ALIGN.CENTER
    else:
        sz = Inches(2.4) if state_index == 0 else Inches(1.35)
        top_y = Inches(1.8) if state_index == 0 else Inches(2.6)
        f1 = slide.shapes.add_picture(IMG_HIBISCUS, x[0] + (w[0] - sz)/2, top_y, sz, sz)
        f1.name = "!!flower_1"

        lbl1 = slide.shapes.add_textbox(x[0], top_y + sz + Inches(0.3), w[0], Inches(1.2))
        lbl1.name = "!!label_1"
        lbl1.text_frame.paragraphs[0].text = "Hibiscus"
        lbl1.text_frame.paragraphs[0].font.name = "Georgia"
        lbl1.text_frame.paragraphs[0].font.size = Pt(18) if state_index == 0 else Pt(13)
        lbl1.text_frame.paragraphs[0].font.italic = True
        lbl1.text_frame.paragraphs[0].font.color.rgb = ACCENT_1
        lbl1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        p_ar = lbl1.text_frame.add_paragraph()
        p_ar.text = "الكركديه"
        p_ar.font.name = "Calibri"
        p_ar.font.size = Pt(14) if state_index == 0 else Pt(11)
        p_ar.font.color.rgb = TEXT_1
        p_ar.alignment = PP_ALIGN.CENTER

    # --- بطاقة 2 ---
    if state_index == 2:
        f2 = slide.shapes.add_picture(IMG_LILY, x[1] + Inches(0.4), Inches(1.2), Inches(4.7), Inches(4.7))
        f2.name = "!!flower_2"

        num2 = slide.shapes.add_textbox(x[1] + Inches(5.0), Inches(0.8), Inches(3.5), Inches(1.8))
        num2.name = "!!num_2"
        p = num2.text_frame.paragraphs[0]
        p.text = "2"
        p.font.name = "Segoe UI"
        p.font.size = Pt(110)
        p.font.bold = True
        p.font.color.rgb = ACCENT_2

        t2 = slide.shapes.add_textbox(x[1] + Inches(4.7), Inches(2.8), Inches(4.2), Inches(3.6))
        t2.name = "!!text_2"
        tf = t2.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = "أيقونة الأناقة والنعومة والجمال الهادئ، وتُعبر عن المشاعر الصادقة والمودة الخالصة."
        p1.font.name = "Calibri"
        p1.font.size = Pt(18)
        p1.font.color.rgb = TEXT_2
        p1.alignment = PP_ALIGN.RIGHT
        p1.space_after = Pt(18)
        p2 = tf.add_paragraph()
        p2.text = "بحضورها الملكي وعطرها الرقيق، تضفي لمسة ساحرة تلائم أرقى مناسبات التهنئة والمحبة."
        p2.font.name = "Calibri"
        p2.font.size = Pt(18)
        p2.font.color.rgb = TEXT_2
        p2.alignment = PP_ALIGN.RIGHT

        lbl2 = slide.shapes.add_textbox(x[1] + Inches(0.7), Inches(5.9), Inches(4.1), Inches(1.2))
        lbl2.name = "!!label_2"
        lbl2.text_frame.paragraphs[0].text = "Pink Lily"
        lbl2.text_frame.paragraphs[0].font.name = "Georgia"
        lbl2.text_frame.paragraphs[0].font.size = Pt(24)
        lbl2.text_frame.paragraphs[0].font.italic = True
        lbl2.text_frame.paragraphs[0].font.color.rgb = ACCENT_2
        lbl2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        p_ar = lbl2.text_frame.add_paragraph()
        p_ar.text = "الزنبق الوردي"
        p_ar.font.name = "Calibri"
        p_ar.font.size = Pt(16)
        p_ar.font.color.rgb = TEXT_2
        p_ar.alignment = PP_ALIGN.CENTER
    else:
        sz = Inches(2.4) if state_index == 0 else Inches(1.35)
        top_y = Inches(1.8) if state_index == 0 else Inches(2.6)
        f2 = slide.shapes.add_picture(IMG_LILY, x[1] + (w[1] - sz)/2, top_y, sz, sz)
        f2.name = "!!flower_2"

        lbl2 = slide.shapes.add_textbox(x[1], top_y + sz + Inches(0.3), w[1], Inches(1.2))
        lbl2.name = "!!label_2"
        lbl2.text_frame.paragraphs[0].text = "Pink Lily"
        lbl2.text_frame.paragraphs[0].font.name = "Georgia"
        lbl2.text_frame.paragraphs[0].font.size = Pt(18) if state_index == 0 else Pt(13)
        lbl2.text_frame.paragraphs[0].font.italic = True
        lbl2.text_frame.paragraphs[0].font.color.rgb = ACCENT_2
        lbl2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        p_ar = lbl2.text_frame.add_paragraph()
        p_ar.text = "الزنبق"
        p_ar.font.name = "Calibri"
        p_ar.font.size = Pt(14) if state_index == 0 else Pt(11)
        p_ar.font.color.rgb = TEXT_2
        p_ar.alignment = PP_ALIGN.CENTER

    # --- بطاقة 3 ---
    if state_index == 3:
        f3 = slide.shapes.add_picture(IMG_PEONY, x[2] + Inches(0.4), Inches(1.2), Inches(4.7), Inches(4.7))
        f3.name = "!!flower_3"

        num3 = slide.shapes.add_textbox(x[2] + Inches(5.0), Inches(0.8), Inches(3.5), Inches(1.8))
        num3.name = "!!num_3"
        p = num3.text_frame.paragraphs[0]
        p.text = "3"
        p.font.name = "Segoe UI"
        p.font.size = Pt(110)
        p.font.bold = True
        p.font.color.rgb = ACCENT_3

        t3 = slide.shapes.add_textbox(x[2] + Inches(4.7), Inches(2.8), Inches(4.2), Inches(3.6))
        t3.name = "!!text_3"
        tf = t3.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = "زهرة نادرة تجسد معاني السكينة والعمق والثقة، وتلفت الأنظار بندرتها وسحرها الفريد."
        p1.font.name = "Calibri"
        p1.font.size = Pt(18)
        p1.font.color.rgb = TEXT_3
        p1.alignment = PP_ALIGN.RIGHT
        p1.space_after = Pt(18)
        p2 = tf.add_paragraph()
        p2.text = "تأسر القلوب بتموجات بتلاتها المخملية العميقة لتعكس فخامة عصرية لا مثيل لها."
        p2.font.name = "Calibri"
        p2.font.size = Pt(18)
        p2.font.color.rgb = TEXT_3
        p2.alignment = PP_ALIGN.RIGHT

        lbl3 = slide.shapes.add_textbox(x[2] + Inches(0.7), Inches(5.9), Inches(4.1), Inches(1.2))
        lbl3.name = "!!label_3"
        lbl3.text_frame.paragraphs[0].text = "Blue Peony"
        lbl3.text_frame.paragraphs[0].font.name = "Georgia"
        lbl3.text_frame.paragraphs[0].font.size = Pt(24)
        lbl3.text_frame.paragraphs[0].font.italic = True
        lbl3.text_frame.paragraphs[0].font.color.rgb = ACCENT_3
        lbl3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        p_ar = lbl3.text_frame.add_paragraph()
        p_ar.text = "الفاوانيا الزرقاء"
        p_ar.font.name = "Calibri"
        p_ar.font.size = Pt(16)
        p_ar.font.color.rgb = TEXT_3
        p_ar.alignment = PP_ALIGN.CENTER
    else:
        sz = Inches(2.4) if state_index == 0 else Inches(1.35)
        top_y = Inches(1.8) if state_index == 0 else Inches(2.6)
        f3 = slide.shapes.add_picture(IMG_PEONY, x[2] + (w[2] - sz)/2, top_y, sz, sz)
        f3.name = "!!flower_3"

        lbl3 = slide.shapes.add_textbox(x[2], top_y + sz + Inches(0.3), w[2], Inches(1.2))
        lbl3.name = "!!label_3"
        lbl3.text_frame.paragraphs[0].text = "Blue Peony"
        lbl3.text_frame.paragraphs[0].font.name = "Georgia"
        lbl3.text_frame.paragraphs[0].font.size = Pt(18) if state_index == 0 else Pt(13)
        lbl3.text_frame.paragraphs[0].font.italic = True
        lbl3.text_frame.paragraphs[0].font.color.rgb = ACCENT_3
        lbl3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        p_ar = lbl3.text_frame.add_paragraph()
        p_ar.text = "الفاوانيا"
        p_ar.font.name = "Calibri"
        p_ar.font.size = Pt(14) if state_index == 0 else Pt(11)
        p_ar.font.color.rgb = TEXT_3
        p_ar.alignment = PP_ALIGN.CENTER

# توليد شرائح PowerPoint
create_powerpoint_slide(0) # نظرة عامة
create_powerpoint_slide(1) # تمدد البطاقة 1
create_powerpoint_slide(2) # تمدد البطاقة 2
create_powerpoint_slide(3) # تمدد البطاقة 3
create_powerpoint_slide(0) # العودة لنظرة عامة

pptx_out = os.path.join(OUTPUT_DIR, "reconstructed_presentation.pptx")
prs.save(pptx_out)
print(f"[1/2] تم حفظ عرض الباوربوينت: {pptx_out}")

# -----------------------------------------------------------------------------
# 5. محرك الاستيفاء الرياضي (Frame-by-Frame Linear Interpolation Engine)
# لتوليد فيديو الحركة بدقة 60 إطاراً في الثانية (Smooth 60 FPS Video)
# -----------------------------------------------------------------------------
try:
    import cv2
    import numpy as np

    WIDTH, HEIGHT = 1920, 1080
    FPS = 60
    TRANSITION_FRAMES = 50   # مدة الحركة الانسيابية (~0.8 ثانية)
    HOLD_FRAMES = 60         # مدة التوقف عند كل بطاقة (1 ثانية)

    # دالة الاستيفاء التكعيبي لتنعيم الحركة (Smoothstep Interpolation)
    def smoothstep(t):
        return t * t * (3 - 2 * t)

    # الحالات الرياضية لكل شريحة (x, width)
    states = [
        # State 0: متساوية
        [(0, 640), (640, 640), (1280, 640)],
        # State 1: الأولى متمددة
        [(0, 1316), (1316, 302), (1618, 302)],
        # State 2: الثانية متمددة
        [(0, 302), (302, 1316), (1618, 302)],
        # State 3: الثالثة متمددة
        [(0, 302), (302, 302), (604, 1316)],
        # State 4: العودة للحالة 0
        [(0, 640), (640, 640), (1280, 640)]
    ]

    bg_colors = [
        (227, 244, 253),  # BGR for Card 1
        (236, 228, 252),  # BGR for Card 2
        (246, 238, 227)   # BGR for Card 3
    ]

    video_out = os.path.join(OUTPUT_DIR, "morph_animation.mp4")
    fourcc = cv2.VideoWriter_fourcc(*'mp4v')
    out = cv2.VideoWriter(video_out, fourcc, FPS, (WIDTH, HEIGHT))

    # تجهيز صور الورود بصيغة ألفا
    raw_imgs = [
        Image.open(IMG_HIBISCUS).convert("RGBA"),
        Image.open(IMG_LILY).convert("RGBA"),
        Image.open(IMG_PEONY).convert("RGBA")
    ]

    print("[2/2] جاري تصيير فيديو الحركة بالاستيفاء الخطي (Interpolation)...")

    for s in range(len(states) - 1):
        start_state = states[s]
        end_state = states[s + 1]

        # 1. إطارات التوقف (Hold)
        for _ in range(HOLD_FRAMES):
            frame = np.ones((HEIGHT, WIDTH, 3), dtype=np.uint8) * 255
            for i in range(3):
                x, w = start_state[i]
                cv2.rectangle(frame, (int(x), 0), (int(x + w), HEIGHT), bg_colors[i], -1)
            out.write(frame)

        # 2. إطارات التحويل التدريجي (Morphing Frames عبر الاستيفاء)
        for f in range(TRANSITION_FRAMES):
            progress = (f + 1) / TRANSITION_FRAMES
            t = smoothstep(progress)  # تطبيق الاستيفاء الرياضي

            frame = np.ones((HEIGHT, WIDTH, 3), dtype=np.uint8) * 255
            for i in range(3):
                # معادلة الاستيفاء الخطي المطبقة:
                # current = start + (end - start) * t
                curr_x = start_state[i][0] + (end_state[i][0] - start_state[i][0]) * t
                curr_w = start_state[i][1] + (end_state[i][1] - start_state[i][1]) * t
                cv2.rectangle(frame, (int(curr_x), 0), (int(curr_x + curr_w), HEIGHT), bg_colors[i], -1)

            out.write(frame)

    out.release()
    print(f"[نجاح]: تم إنشاء فيديو الحركة السلسة بنجاح: {video_out}")

except ImportError:
    print("[تنبيه]: لتوليد فيديو MP4 للحركة تلقائياً، قم بتثبيت opencv عبر: pip install opencv-python-headless")
